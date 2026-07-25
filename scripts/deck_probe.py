#!/usr/bin/env python3
"""deck_probe.py — deterministic *reusable-form* probe for the assimilate-slides
skill's INGEST step.

Why this exists: without a fixed tool, an agent assimilating a deck improvises a
parser and `pip install`s whatever it guesses (pymupdf, pdf2image, cairosvg,
unoconv, …) — non-deterministic and against the repo's no-new-dependency rule.
This script does the extraction with **only already-present deps**
(`python-pptx`, `lxml`, stdlib) and, for formats no bundled lib handles
(PDF / images / SVG), prints guidance to use the harness's own file viewer
rather than installing anything.

It reports *style/structure*, not content: palette, fonts, per-slide shape
composition, dense-diagram slides, and image-heavy discard candidates. Short
text labels are printed only behind `--labels` (for icon-concept spotting).

Usage:
    python3 scripts/deck_probe.py <file|dir> [--labels]

Output goes to stdout — redirect into the gitignored `.context/` scratch; never
commit it. No network, no install, no writes.
"""
from __future__ import annotations

import collections
import re
import shutil
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path

SUPPORTED_PPTX = {".pptx"}
SUPPORTED_WEB = {".html", ".htm", ".css"}
VISUAL = {".png", ".jpg", ".jpeg", ".svg", ".webp", ".gif"}


def probe_pptx(path: Path, show_labels: bool) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        print("python-pptx not installed. It is a declared repo dep: "
              "`pip install python-pptx lxml`. Do NOT substitute another library.",
              file=sys.stderr)
        raise SystemExit(2)

    prs = Presentation(str(path))
    W, H = prs.slide_width, prs.slide_height
    print(f"SOURCE type=pptx slides={len(prs.slides)} "
          f"size={Emu(W).inches:.2f}x{Emu(H).inches:.2f}in")  # type: ignore[arg-type]

    fonts: collections.Counter = collections.Counter()
    fill_colors: collections.Counter = collections.Counter()
    font_colors: collections.Counter = collections.Counter()
    dense, discard = [], []

    for i, slide in enumerate(prs.slides, 1):
        n_pic = n_txt = n_tbl = n_grp = 0
        pic_area = chars = 0
        auto = free = line = 0
        for shp in slide.shapes:
            st = str(shp.shape_type or "")
            if "PICTURE" in st:
                n_pic += 1
                try:
                    pic_area += (shp.width or 0) * (shp.height or 0)
                except Exception:
                    pass
            if "AUTO_SHAPE" in st:
                auto += 1
            if "FREEFORM" in st:
                free += 1
            if st.startswith("LINE") or "CONNECTOR" in st:
                line += 1
            if "GROUP" in st:
                n_grp += 1
            if getattr(shp, "has_table", False) and shp.has_table:
                n_tbl += 1
            if shp.has_text_frame:
                n_txt += 1
                chars += len(shp.text_frame.text)
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name:
                            fonts[r.font.name] += 1
                        try:
                            c = r.font.color
                            if c and c.type is not None and c.rgb is not None:
                                font_colors[str(c.rgb)] += 1
                        except Exception:
                            pass
            try:
                f = shp.fill
                if f.type is not None and "SOLID" in str(f.type) and f.fore_color and f.fore_color.rgb:
                    fill_colors[str(f.fore_color.rgb)] += 1
            except Exception:
                pass
        frac = (pic_area / (W * H)) if (W and H) else 0
        diagram_shapes = auto + free + line
        flags = []
        if frac > 0.55 and chars < 120:
            flags.append("IMG_HEAVY→discard")
            discard.append(i)
        if diagram_shapes >= 20:
            flags.append("DENSE_DIAGRAM")
            dense.append(i)
        print(f"S{i:02d} layout={slide.slide_layout.name!r} pic={n_pic}({frac:.0%}) "
              f"txt={n_txt} chars={chars} tbl={n_tbl} grp={n_grp} "
              f"auto/free/line={auto}/{free}/{line} {' '.join(flags)}")

    print("\nTOP_FONTS", fonts.most_common(10))
    print("FOCUS/FONT_COLORS", font_colors.most_common(12))
    print("FILL_COLORS", fill_colors.most_common(12))
    print("DENSE_DIAGRAM_SLIDES (architecture-canvas candidates):", dense)
    print("DISCARD_CANDIDATES (image-only headers/dividers):", discard)

    if show_labels:
        seen = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if shp.has_text_frame:
                    t = shp.text_frame.text.strip().replace("\n", " ")
                    if 0 < len(t) <= 40 and t not in seen:
                        seen.append(t)
        print("\nSHORT_LABELS (for icon-concept spotting — scrub before any commit):")
        print(" | ".join(seen[:60]))


# ── Brand-shell extraction ────────────────────────────────────────────────


@dataclass
class BrandShell:
    """Raw brand signals extracted from a PPTX — no base-style selection here."""
    detected_mode: str = "neutral"   # "dark" | "light" | "neutral"
    bg_primary: str = ""
    bg_secondary: str = ""
    text_primary: str = ""
    text_secondary: str = ""
    accent_1: str = ""
    accent_2: str = ""
    accent_3: str = ""
    accent_4: str = ""
    footer_text: str = ""
    fill_colors: list[str] = field(default_factory=list)
    sparse: bool = False


def _relative_luminance(hex_color: str) -> float:
    """WCAG 2.1 relative luminance from a 6-char hex string (with or without #)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255

    def _lin(c: float) -> float:
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4

    return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)


def _mode_from_luminance(L: float) -> str:
    if L <= 0.18:
        return "dark"
    if L >= 0.40:
        return "light"
    return "neutral"


def probe_pptx_brand(path: Path, footer_text: bool) -> BrandShell:
    """Extract brand color signals from a PPTX without selecting a base style.

    Reads only scalar fill/font colors — no PICTURE pixel access.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("python-pptx not installed.", file=sys.stderr)
        raise SystemExit(2)

    prs = Presentation(str(path))
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    header_threshold = slide_height * (80 / 720)   # ~11% of standard 720pt height
    footer_threshold = slide_height * 0.9

    fill_counter: collections.Counter = collections.Counter()
    header_colors: list[str] = []
    text_colors: list[str] = []
    footer_candidates: list[str] = []
    bg_hex: str = ""

    for slide in prs.slides:
        # Slide background — canonical luminance source
        bg = slide.background.fill
        try:
            if bg.type is not None and "SOLID" in str(bg.type):
                rgb = bg.fore_color.rgb
                bg_hex = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
        except Exception:
            pass

        for shp in slide.shapes:
            shape_type = str(shp.shape_type or "")
            if "PICTURE" in shape_type:
                continue  # never access pixel content

            # Solid fills on AUTO_SHAPE and FREEFORM
            if "AUTO_SHAPE" in shape_type or "FREEFORM" in shape_type:
                try:
                    f = shp.fill
                    if f.type is not None and "SOLID" in str(f.type) and f.fore_color and f.fore_color.rgb:
                        rgb = f.fore_color.rgb
                        hex_color = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                        fill_counter[hex_color] += 1
                        try:
                            top = shp.top
                            if top is not None and top < header_threshold:
                                header_colors.append(hex_color)
                        except Exception:
                            pass
                except Exception:
                    pass

            # Text run font colors + footer text
            if shp.has_text_frame:
                try:
                    for para in shp.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                c = run.font.color
                                if c and c.type is not None and c.rgb is not None:
                                    rgb = c.rgb
                                    text_colors.append(f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
                            except Exception:
                                pass
                    if footer_text:
                        try:
                            top = shp.top
                            text = shp.text_frame.text.strip()
                            if (top is not None and top >= footer_threshold
                                    and 0 < len(text) <= 60):
                                footer_candidates.append(text)
                        except Exception:
                            pass
                except Exception:
                    pass

    # Determine background and luminance
    if bg_hex:
        bg_lum = _relative_luminance(bg_hex)
        detected_mode = _mode_from_luminance(bg_lum)
    elif fill_counter:
        # Fallback: most-frequent shape fill
        bg_hex = fill_counter.most_common(1)[0][0]
        bg_lum = _relative_luminance(bg_hex)
        detected_mode = _mode_from_luminance(bg_lum)
    else:
        detected_mode = "neutral"

    distinct_fills = list(fill_counter.keys())
    sparse = len(distinct_fills) < 3

    # Map collected colors into BrandShell slots
    # bg_secondary: second most common fill (skip bg_primary)
    fill_by_freq = [h for h, _ in fill_counter.most_common() if h != bg_hex]
    bg_secondary = fill_by_freq[0] if fill_by_freq else ""

    # accents: header colors take priority, then remaining fills
    accent_pool: list[str] = []
    seen_accents: set[str] = set()
    for c in header_colors:
        if c not in seen_accents and c != bg_hex:
            accent_pool.append(c)
            seen_accents.add(c)
    for c in fill_by_freq:
        if c not in seen_accents:
            accent_pool.append(c)
            seen_accents.add(c)

    # text colors: pick two most common
    text_counter: collections.Counter = collections.Counter(text_colors)
    top_text = [h for h, _ in text_counter.most_common(2)]

    return BrandShell(
        detected_mode=detected_mode,
        bg_primary=f"#{bg_hex}" if bg_hex else "",
        bg_secondary=f"#{bg_secondary}" if bg_secondary else "",
        text_primary=f"#{top_text[0]}" if len(top_text) > 0 else "",
        text_secondary=f"#{top_text[1]}" if len(top_text) > 1 else "",
        accent_1=f"#{accent_pool[0]}" if len(accent_pool) > 0 else "",
        accent_2=f"#{accent_pool[1]}" if len(accent_pool) > 1 else "",
        accent_3=f"#{accent_pool[2]}" if len(accent_pool) > 2 else "",
        accent_4=f"#{accent_pool[3]}" if len(accent_pool) > 3 else "",
        footer_text=footer_candidates[0] if footer_candidates else "",
        fill_colors=list(fill_counter.keys()),
        sparse=sparse,
    )


def probe_pdf(path: Path, show_labels: bool) -> None:
    """Ingest a source PDF *install-free*: use poppler's pdftotext/pdfinfo when
    already present; otherwise route to the harness viewer. Never installs a
    renderer/parser (no pymupdf/pdf2image/poppler-via-pip)."""
    print(f"SOURCE type=pdf file={path.name}")
    info = shutil.which("pdfinfo")
    totext = shutil.which("pdftotext")
    if info:
        try:
            out = subprocess.run([info, str(path)], capture_output=True, text=True, timeout=30).stdout
            pages = next((l.split(":", 1)[1].strip() for l in out.splitlines() if l.startswith("Pages")), "?")
            print(f"pages={pages} (via poppler pdfinfo)")
        except Exception:
            pass
    if totext:
        try:
            txt = subprocess.run([totext, "-layout", str(path), "-"], capture_output=True,
                                 text=True, timeout=60).stdout
            print(f"text_chars={len(txt)} (via poppler pdftotext)")
            if show_labels:
                labels = [t.strip() for t in txt.splitlines() if 0 < len(t.strip()) <= 40]
                uniq = list(dict.fromkeys(labels))
                print("\nSHORT_LABELS (icon-concept spotting — scrub before any commit):")
                print(" | ".join(uniq[:60]))
        except Exception:
            pass
    else:
        print("poppler not installed — read this PDF with the harness's own file viewer "
              "(it renders PDF pages); do NOT `pip install` pymupdf/pdf2image/weasyprint or "
              "`apt install poppler`. Extract palette/layout/icon-ideas by eye.")


def probe_web(path: Path) -> None:
    text = path.read_text(encoding="utf-8", errors="replace")
    print(f"SOURCE type=web/{path.suffix.lstrip('.')} bytes={len(text)}")
    root_vars = re.findall(r"(--[\w-]+)\s*:\s*([^;]+);", text)
    fonts = re.findall(r"font-family\s*:\s*([^;{}]+)[;}]", text, re.I)
    hexes = collections.Counter(re.findall(r"#[0-9a-fA-F]{3,8}\b", text))
    print("\nCSS_ROOT_VARS:")
    for k, v in root_vars[:40]:
        print(f"  {k}: {v.strip()[:60]}")
    print("\nFONT_STACKS:")
    for f in dict.fromkeys(s.strip()[:70] for s in fonts):
        print(f"  {f}")
    print("\nTOP_HEX_COLORS", hexes.most_common(15))


def main() -> int:
    show_brand = "--brand" in sys.argv
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    show_labels = "--labels" in sys.argv
    if not args:
        print(__doc__)
        return 2
    target = Path(args[0]).expanduser()
    if not target.exists():
        print(f"not found: {target}", file=sys.stderr)
        return 1

    if show_brand:
        # Brand diagnostic: raw extraction only, no base-style selection
        if target.suffix.lower() != ".pptx":
            print(f"--brand requires a .pptx file, got: {target}", file=sys.stderr)
            return 1
        shell = probe_pptx_brand(target, footer_text=True)
        print("BRAND_SHELL")
        print(f"  detected_mode:   {shell.detected_mode}")
        print(f"  bg_primary:      {shell.bg_primary}")
        print(f"  bg_secondary:    {shell.bg_secondary}")
        print(f"  text_primary:    {shell.text_primary}")
        print(f"  text_secondary:  {shell.text_secondary}")
        print(f"  accent_1:        {shell.accent_1}")
        print(f"  accent_2:        {shell.accent_2}")
        print(f"  accent_3:        {shell.accent_3}")
        print(f"  accent_4:        {shell.accent_4}")
        print(f"  fill_colors:     {shell.fill_colors}")
        print(f"  sparse:          {shell.sparse}")
        if shell.footer_text:
            print(f"  footer_text:     {shell.footer_text!r}")
        return 0

    paths = sorted(target.rglob("*")) if target.is_dir() else [target]
    handled = False
    for p in paths:
        if p.is_dir():
            continue
        ext = p.suffix.lower()
        if ext in SUPPORTED_PPTX:
            print("=" * 72)
            probe_pptx(p, show_labels)
            handled = True
        elif ext in SUPPORTED_WEB:
            print("=" * 72)
            probe_web(p)
            handled = True
        elif ext in VISUAL:
            print("=" * 72)
            print(f"SOURCE type={ext.lstrip('.')} — visual/binary. Read it with the "
                  f"harness's own file viewer (it renders PDFs/images); do NOT install "
                  f"a renderer or SVG/PDF parser. Extract palette/layout/icon-ideas by eye.")
            handled = True
    if not handled:
        print("no supported files found (looked for: .pptx, .html/.css; "
              ".pdf/.png/.jpg/.svg are read by the harness viewer, not this script).")
    return 0


if __name__ == "__main__":
    sys.exit(main())
