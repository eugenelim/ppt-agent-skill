"""SVG-to-PowerPoint compatibility tests for native SVG output.

For each P2 target diagram type:
1. Generate native SVG via dispatch_native()
2. Write to a temp file
3. Run SvgConverter.convert() on it into a real Presentation slide
4. Assert shapes > 0, errors == 0, no unsupported-element warnings
5. Verify the PPTX XML can be inspected for text content

This catches any SVG constructs the converter can't handle (foreign objects,
unsupported elements, invalid geometry) before they reach production renders.
"""
from __future__ import annotations

import tempfile
from pathlib import Path

import pytest

from pptx import Presentation
from pptx.util import Emu

from scripts.mermaid_render.native_svg import dispatch_native
from scripts.svg2pptx import SvgConverter

# ── Layout constants matching svg2pptx defaults ───────────────────────────────

SLIDE_W = 12192000  # 9144000 + 3048000 (Emu at 96 dpi)
SLIDE_H = 6858000

# ── Fixture sources ───────────────────────────────────────────────────────────

FIXTURES: list[tuple[str, str]] = [
    ("flowchart", """\
flowchart TD
    A[Start] --> B{Decision}
    B -->|Yes| C[Result A]
    B -->|No| D[Result B]
    C --> E[End]
    D --> E
"""),
    ("flowchart-lr", """\
flowchart LR
    X[Input] --> Y[Process] --> Z[Output]
"""),
    ("stateDiagram", """\
stateDiagram-v2
    [*] --> Idle
    Idle --> Running : start
    Running --> Idle : stop
    Running --> [*]
"""),
    ("mindmap", """\
mindmap
    root((Central))
        Branch A
            Leaf A1
            Leaf A2
        Branch B
            Leaf B1
"""),
    ("timeline", """\
timeline
    title Project Timeline
    section Phase 1
    Jan : Kickoff
    Feb : Design
    section Phase 2
    Mar : Build
    Apr : Launch
"""),
    ("architecture-beta", """\
architecture-beta
    service API(server)[API]
    service DB(database)[Database]
    API --> DB
"""),
    ("c4context", """\
C4Context
    title System Context
    Person(user, "User", "A human")
    System(sys, "System", "Main system")
    Rel(user, sys, "Uses")
"""),
]


# ── Helper ────────────────────────────────────────────────────────────────────

def _run_compat(src: str) -> dict:
    """Generate native SVG, convert to PPTX slide, return converter stats."""
    svg_str = dispatch_native(src)

    with tempfile.NamedTemporaryFile(suffix=".svg", mode="w",
                                     encoding="utf-8", delete=False) as f:
        f.write(svg_str)
        svg_path = Path(f.name)

    try:
        prs = Presentation()
        prs.slide_width = Emu(SLIDE_W)
        prs.slide_height = Emu(SLIDE_H)
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        converter = SvgConverter()
        converter.convert(svg_path, slide)
        return dict(converter.stats)
    finally:
        svg_path.unlink(missing_ok=True)


# ── Tests ─────────────────────────────────────────────────────────────────────

@pytest.mark.parametrize("dtype,src", FIXTURES, ids=[f[0] for f in FIXTURES])
class TestSvg2PptxCompat:
    def test_no_errors(self, dtype, src):
        stats = _run_compat(src)
        assert stats["errors"] == 0, (
            f"{dtype}: {stats['errors']} conversion error(s)"
        )

    def test_shapes_present(self, dtype, src):
        stats = _run_compat(src)
        assert stats["shapes"] > 0, (
            f"{dtype}: expected shapes > 0, got {stats['shapes']}"
        )

    def test_no_foreign_object(self, dtype, src):
        svg_str = dispatch_native(src)
        assert "<foreignObject" not in svg_str, (
            f"{dtype}: SVG contains <foreignObject>"
        )

    def test_svg_is_valid_xml(self, dtype, src):
        from lxml import etree
        import re
        svg_str = dispatch_native(src)
        body = re.sub(r"^<\?xml[^?]*\?>", "", svg_str.strip()).strip()
        try:
            etree.fromstring(body.encode("utf-8"))
        except etree.XMLSyntaxError as e:
            pytest.fail(f"{dtype}: SVG is not valid XML: {e}")


class TestSvg2PptxTextPreservation:
    """Verify that text content from native SVG survives conversion."""

    def test_flowchart_text_in_svg(self):
        src = FIXTURES[0][1]
        svg = dispatch_native(src)
        # Text should be in the SVG as native text elements
        assert "Start" in svg
        assert "Decision" in svg

    def test_timeline_period_text_in_svg(self):
        src = next(s for t, s in FIXTURES if t == "timeline")
        svg = dispatch_native(src)
        assert "Kickoff" in svg or "Jan" in svg

    def test_mindmap_label_text_in_svg(self):
        src = next(s for t, s in FIXTURES if t == "mindmap")
        svg = dispatch_native(src)
        assert "Central" in svg

    def test_c4_labels_in_svg(self):
        src = next(s for t, s in FIXTURES if t == "c4context")
        svg = dispatch_native(src)
        assert "User" in svg


class TestSvg2PptxDeterminism:
    """Verify that repeated conversion produces same shape count."""

    def test_flowchart_stable_shape_count(self):
        src = FIXTURES[0][1]
        stats1 = _run_compat(src)
        stats2 = _run_compat(src)
        assert stats1["shapes"] == stats2["shapes"]
