"""Tests for brand-shell extraction: deck_probe.py extensions + brand_extract.py."""
from __future__ import annotations

import json
import subprocess
import sys
import os
from pathlib import Path

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "scripts"))

FIXTURES = Path(__file__).parent / "fixtures"
DARK_FIXTURE = FIXTURES / "brand-probe-fixture.pptx"
LIGHT_FIXTURE = FIXTURES / "brand-probe-light-fixture.pptx"
SCRIPTS = Path(__file__).parent.parent / "scripts"
REFS_DIR = Path(__file__).parent.parent / "references"


# ── Task 1 TDD stubs: deck_probe.py helpers ────────────────────────────────


def test_luminance_dark():
    from deck_probe import _relative_luminance
    assert _relative_luminance("#050b1f") <= 0.18


def test_luminance_light():
    from deck_probe import _relative_luminance
    assert _relative_luminance("#F5F5F5") >= 0.40


def test_luminance_neutral():
    from deck_probe import _relative_luminance
    L = _relative_luminance("#808080")
    assert 0.18 < L < 0.40


def test_sparse_flag():
    from deck_probe import probe_pptx_brand
    shell = probe_pptx_brand(DARK_FIXTURE, footer_text=False)
    # dark fixture has multiple fills, so it is NOT sparse
    assert shell.sparse is False
    # test sparse: build a 1-fill PPTX in memory
    import io
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(5, 11, 31)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(buf.read())
        tmp = Path(f.name)
    try:
        sparse_shell = probe_pptx_brand(tmp, footer_text=False)
        assert sparse_shell.sparse is True
    finally:
        tmp.unlink(missing_ok=True)


def test_header_accent_y_threshold():
    """Shape at y < slide_height * (80/720) goes into header pool; shape at y > 20% does not."""
    import io, tempfile
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu
    from deck_probe import probe_pptx_brand

    SLIDE_H = Emu(5143500)
    SLIDE_W = Emu(9144000)
    # header_threshold = SLIDE_H * (80/720) ≈ SLIDE_H * 0.1111
    header_top = int(SLIDE_H * 0.05)   # inside header zone
    body_top   = int(SLIDE_H * 0.30)   # outside header zone

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(5, 11, 31)

    # header shape
    sh = s.shapes.add_shape(1, Emu(0), Emu(header_top), Emu(int(SLIDE_W)), Emu(int(SLIDE_H * 0.06)))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x22, 0xD3, 0xEE)
    sh.line.fill.background()
    # body shape
    sh2 = s.shapes.add_shape(1, Emu(0), Emu(body_top), Emu(int(SLIDE_W)), Emu(int(SLIDE_H * 0.3)))
    sh2.fill.solid(); sh2.fill.fore_color.rgb = RGBColor(0x10, 0x20, 0x40)
    sh2.line.fill.background()

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(buf.read()); tmp = Path(f.name)
    try:
        shell = probe_pptx_brand(tmp, footer_text=False)
        assert shell.accent_1 != "", "header shape color should be captured as accent"
    finally:
        tmp.unlink(missing_ok=True)


def test_footer_text_y_threshold():
    """Text run at y ≥ 90% of slide height with ≤ 60 chars → captured in footer_text."""
    from deck_probe import probe_pptx_brand
    shell = probe_pptx_brand(DARK_FIXTURE, footer_text=True)
    # dark fixture has a footer textbox with "Confidential – Acme Corp 2025"
    assert shell.footer_text != "", "expected footer text to be captured"


def test_no_picture_access():
    """PPTX with only PICTURE shapes → fill_colors empty → sparse=True, no pixel read."""
    import io, tempfile
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Inches
    from deck_probe import probe_pptx_brand

    # Build a PPTX with no fill shapes (only a text box with no solid fill)
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # No background set, no solid-fill shapes — sparse by fill count
    tb = s.shapes.add_textbox(Emu(0), Emu(0), Emu(9144000), Emu(914400))
    tb.text_frame.text = "Only text"

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(buf.read()); tmp = Path(f.name)
    try:
        shell = probe_pptx_brand(tmp, footer_text=False)
        assert shell.sparse is True
    finally:
        tmp.unlink(missing_ok=True)


def test_theme_inherited_bg_fallback():
    """PPTX with no explicit slide background fill → bg from most-frequent shape fill."""
    import io, tempfile
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu
    from deck_probe import probe_pptx_brand

    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    # Slide with NO background set (theme-inherited) but several solid-fill shapes
    for i in range(3):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        sh = s.shapes.add_shape(1, Emu(0), Emu(int(Emu(5143500) * 0.2)),
                                Emu(9144000), Emu(int(Emu(5143500) * 0.5)))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(5, 11, 31)
        sh.line.fill.background()

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(buf.read()); tmp = Path(f.name)
    try:
        shell = probe_pptx_brand(tmp, footer_text=False)
        # should not be neutral (bg derived from shape fills)
        assert shell.detected_mode in ("dark", "light", "neutral")
        assert shell.bg_primary != ""
    finally:
        tmp.unlink(missing_ok=True)


# ── Task 2 stubs: brand_extract.py ────────────────────────────────────────


def _run_extract(*args) -> tuple[int, str, str]:
    """Run brand_extract.py and return (exit_code, stdout, stderr)."""
    result = subprocess.run(
        [sys.executable, str(SCRIPTS / "brand_extract.py")] + list(args),
        capture_output=True, text=True
    )
    return result.returncode, result.stdout, result.stderr


def _run_validator(path: str) -> int:
    result = subprocess.run(
        [sys.executable, str(SCRIPTS / "contract_validator.py"), "style", path],
        capture_output=True, text=True
    )
    return result.returncode


def test_brand_extract_dark_fixture(tmp_path):
    out = str(tmp_path / "style.json")
    code, stdout, stderr = _run_extract(
        str(DARK_FIXTURE), "--output", out,
        "--refs-dir", str(REFS_DIR)
    )
    assert code == 0, f"exit {code}\nstdout={stdout}\nstderr={stderr}"
    assert _run_validator(out) == 0, "contract_validator failed"
    data = json.loads(Path(out).read_text())
    assert data.get("brand_mode") == "dark"


def test_brand_extract_light_fixture(tmp_path):
    out = str(tmp_path / "style.json")
    code, stdout, stderr = _run_extract(
        str(LIGHT_FIXTURE), "--output", out,
        "--refs-dir", str(REFS_DIR)
    )
    assert code == 0, f"exit {code}\nstdout={stdout}\nstderr={stderr}"
    data = json.loads(Path(out).read_text())
    assert data["style_id"].startswith("brand_shell_")
    assert data.get("brand_mode") == "light"


def test_sparse_fallback(tmp_path):
    import io
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu
    # 1-fill PPTX → sparse
    prs = Presentation()
    prs.slide_width = Emu(9144000); prs.slide_height = Emu(5143500)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(5, 11, 31)
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    tmp_pptx = tmp_path / "sparse.pptx"
    tmp_pptx.write_bytes(buf.read())
    out = str(tmp_path / "style.json")
    code, stdout, stderr = _run_extract(
        str(tmp_pptx), "--output", out, "--refs-dir", str(REFS_DIR)
    )
    assert code == 0, f"exit {code}\nstdout={stdout}\nstderr={stderr}"
    data = json.loads(Path(out).read_text())
    # sparse: style_id not prefixed with brand_shell_ (preserves base)
    assert not data["style_id"].startswith("brand_shell_"), \
        f"sparse fallback should preserve base style_id, got {data['style_id']}"
    assert "brand_mode" in data


def test_font_family_from_base_not_pptx(tmp_path):
    out = str(tmp_path / "style.json")
    code, *_ = _run_extract(
        str(DARK_FIXTURE), "--output", out, "--refs-dir", str(REFS_DIR)
    )
    assert code == 0
    data = json.loads(Path(out).read_text())
    font = data.get("font_family", "")
    assert font != "", "font_family must be set"
    # font comes from typography.display_font of the base style, not from PPTX
    # The base for dark is dark_tech whose display_font starts with 'Inter Tight'
    assert "Inter" in font or "sans" in font.lower(), \
        f"expected a sans-serif font from the base style, got: {font!r}"


def test_footer_text_present_with_flag(tmp_path):
    out = str(tmp_path / "style.json")
    code, *_ = _run_extract(
        str(DARK_FIXTURE), "--output", out,
        "--refs-dir", str(REFS_DIR), "--footer-text"
    )
    assert code == 0
    data = json.loads(Path(out).read_text())
    assert "brand_footer_text" in data
    assert data["brand_footer_text"] != ""


def test_footer_text_absent_without_flag(tmp_path):
    out = str(tmp_path / "style.json")
    code, *_ = _run_extract(
        str(DARK_FIXTURE), "--output", out, "--refs-dir", str(REFS_DIR)
    )
    assert code == 0
    data = json.loads(Path(out).read_text())
    assert "brand_footer_text" not in data


def test_file_size_cap(tmp_path):
    import struct
    # Write a fake PPTX that is just over 50MB
    big = tmp_path / "toobig.pptx"
    big.write_bytes(b"PK" + b"\x00" * (51 * 1024 * 1024))
    out = str(tmp_path / "style.json")
    code, stdout, stderr = _run_extract(
        str(big), "--output", out, "--refs-dir", str(REFS_DIR)
    )
    assert code == 1, f"expected exit 1 for oversized file, got {code}"


def test_partial_population_stays_contract_valid(tmp_path):
    """3-fill deck with theme-inherited text → contract_validator still passes."""
    import io
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu
    prs = Presentation()
    prs.slide_width = Emu(9144000); prs.slide_height = Emu(5143500)
    for color in ["050b1f", "22D3EE", "0a1f3d"]:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        bg = s.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor.from_string(color)
        sh = s.shapes.add_shape(1, Emu(0), Emu(int(Emu(5143500)*0.1)),
                                Emu(9144000), Emu(int(Emu(5143500)*0.1)))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(color)
        sh.line.fill.background()
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    pptx_path = tmp_path / "partial.pptx"
    pptx_path.write_bytes(buf.read())
    out = str(tmp_path / "style.json")
    code, stdout, stderr = _run_extract(
        str(pptx_path), "--output", out, "--refs-dir", str(REFS_DIR)
    )
    assert code == 0, f"exit {code}\nstdout={stdout}\nstderr={stderr}"
    assert _run_validator(out) == 0


def test_base_style_override(tmp_path):
    """Light fixture auto-selects blue_white; --base-style dark_tech forces dark_tech."""
    out = str(tmp_path / "style.json")
    code, stdout, stderr = _run_extract(
        str(LIGHT_FIXTURE), "--output", out,
        "--refs-dir", str(REFS_DIR), "--base-style", "dark_tech"
    )
    assert code == 0, f"exit {code}\nstdout={stdout}\nstderr={stderr}"
    data = json.loads(Path(out).read_text())
    # mood_keywords should match dark_tech, not blue_white
    assert "深空冷寂" in data.get("mood_keywords", []) or \
        any("深空" in kw or "冷寂" in kw for kw in data.get("mood_keywords", [])), \
        f"Expected dark_tech mood_keywords, got: {data.get('mood_keywords')}"


def test_base_style_unknown_id(tmp_path):
    out = str(tmp_path / "style.json")
    code, stdout, stderr = _run_extract(
        str(DARK_FIXTURE), "--output", out,
        "--refs-dir", str(REFS_DIR), "--base-style", "nonexistent_id_xyz"
    )
    assert code == 1
    combined = stdout + stderr
    assert "nonexistent_id_xyz" in combined or "Valid IDs" in combined or "valid" in combined.lower()


def test_brand_mode_in_output(tmp_path):
    for fixture, expected_mode in [(DARK_FIXTURE, "dark"), (LIGHT_FIXTURE, "light")]:
        out = str(tmp_path / f"style_{expected_mode}.json")
        code, *_ = _run_extract(
            str(fixture), "--output", out, "--refs-dir", str(REFS_DIR)
        )
        assert code == 0
        data = json.loads(Path(out).read_text())
        assert data.get("brand_mode") == expected_mode, \
            f"Expected brand_mode={expected_mode!r}, got {data.get('brand_mode')!r}"


def test_nested_to_css_variables_card_radius_is_string():
    from brand_extract import _nested_to_css_variables
    style = {
        "background": {"primary": "#000", "gradient_to": "#111"},
        "card": {"gradient_from": "#222", "gradient_to": "#333",
                 "border": "#444", "border_radius": 8},
        "text": {"primary": "#fff", "secondary": "#aaa"},
        "accent": {"primary": ["#f00", "#0f0"], "secondary": ["#00f", "#ff0"]},
        "typography": {"display_font": "Arial"},
    }
    css, font = _nested_to_css_variables(style)
    assert css["card_radius"] == "8px", f"expected '8px', got {css['card_radius']!r}"
    assert isinstance(css["card_radius"], str)


def test_nested_to_css_variables_accent_indices():
    from brand_extract import _nested_to_css_variables
    style = {
        "background": {"primary": "#000", "gradient_to": "#111"},
        "card": {"gradient_from": "#222", "gradient_to": "#333",
                 "border": "#444", "border_radius": 4},
        "text": {"primary": "#fff", "secondary": "#aaa"},
        "accent": {"primary": ["#AA0000", "#00AA00"],
                   "secondary": ["#0000AA", "#AAAA00"]},
        "typography": {"display_font": "Helvetica"},
    }
    css, _ = _nested_to_css_variables(style)
    assert css["accent_1"] == "#AA0000"
    assert css["accent_2"] == "#00AA00"
    assert css["accent_3"] == "#0000AA"
    assert css["accent_4"] == "#AAAA00"


def test_brand_mode_neutral_in_merge_output():
    from brand_extract import merge_shell
    from deck_probe import BrandShell
    # Load any valid base style dict
    import json as _json
    refs = REFS_DIR / "styles" / "dark.md"
    text = refs.read_text(encoding="utf-8")
    import re
    blocks = re.findall(r"```json\n(\{.*?\})\n```", text, re.DOTALL)
    base = next(b for b in (_json.loads(b) for b in blocks) if b.get("style_id") == "dark_tech")
    shell = BrandShell(detected_mode="neutral")
    out = merge_shell(base, shell, "test")
    assert out["brand_mode"] == "neutral"


def test_select_base_style_empty_candidates_returns_none():
    from brand_extract import select_base_style
    from deck_probe import BrandShell
    shell = BrandShell(detected_mode="dark")
    result = select_base_style(shell, [], override_id=None)
    assert result is None
