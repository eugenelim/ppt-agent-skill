"""Stage 13 Task C: SVG-to-PowerPoint compatibility pass.

For each PARTIAL/FULL type in NATIVE_RENDERER_REGISTRY, render a minimal
source to native SVG then pass it through scripts/svg2pptx.py SvgConverter.
Asserts no crash and that the embed path returns a valid slide with ≥0 errors.
"""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_ROOT / "scripts"))

from svg2pptx import SvgConverter  # noqa: E402
from mermaid_render.native_svg import dispatch_native  # noqa: E402
from mermaid_render.scene import NATIVE_RENDERER_REGISTRY, NativeParityLevel  # noqa: E402


# Minimal sources per directive — same set as _REGISTRY_SEMANTIC_PARAMS in
# test_native_renderer_capabilities.py, kept in sync by inspection.
_PPTX_COMPAT_SOURCES: dict[str, str] = {
    "flowchart":          "flowchart LR\n  A[Alpha] --> B[Beta]",
    "graph":              "graph TD\n  A --> B",
    "statediagram-v2":    "stateDiagram-v2\n  [*] --> Active\n  Active --> [*]",
    "statediagram":       "stateDiagram\n  [*] --> A\n  A --> [*]",
    "classdiagram":       "classDiagram\n  class Animal {\n    +name : string\n  }",
    "timeline":           "timeline\n  title My Timeline\n  2020 : Launch",
    "mindmap":            "mindmap\n  root((Root))\n    Item A",
    "architecture-beta":  "architecture-beta\n  service svc(internet)[Svc]",
    "c4context":          'C4Context\n  Person(p, "User")\n  System(s, "App")',
    "c4container":        'C4Container\n  Container(c, "API", "Python")',
    "c4component":        'C4Component\n  Component(cp, "Svc", "Python")',
    "sequencediagram":    "sequenceDiagram\n  Alice->>Bob: Hello",
    "erdiagram":          "erDiagram\n  PERSON { string name }",
    "gantt":              "gantt\n  title G\n  section A\n    Task1 :t1, 2024-01-01, 7d",
    "quadrantchart":      "quadrantChart\n  x-axis Low --> High\n  y-axis Low --> High\n  P: [0.5, 0.7]",
    "pie":                'pie\n  title Pets\n  "Dogs" : 386',
    "xychart-beta":       "xychart-beta\n  x-axis [a, b, c]\n  y-axis 0 --> 10\n  bar [5, 3, 8]",
    "block-beta":         "block-beta\n  A B C",
    "packet-beta":        "packet-beta\n  0-7: Source Port",
    "kanban":             "kanban\n  column1\n    item1[Task 1]",
    "journey":            "journey\n  title My day\n  section Go\n    Task: 5: Me",
    "requirementdiagram": "requirementDiagram\n  requirement req1 {\n    id: 1\n    text: Example\n  }",
    "gitgraph":           'gitGraph\n  commit\n  commit id: "second"',
}

_PARTIAL_FULL = [
    d for d, spec in NATIVE_RENDERER_REGISTRY.items()
    if spec.parity in (NativeParityLevel.FULL, NativeParityLevel.PARTIAL)
]


def test_pptx_compat_sources_cover_all_partial_full():
    """_PPTX_COMPAT_SOURCES must have an entry for every PARTIAL/FULL registry directive."""
    missing = [d for d in _PARTIAL_FULL if d not in _PPTX_COMPAT_SOURCES]
    assert not missing, (
        f"Add entries to _PPTX_COMPAT_SOURCES for: {missing!r}"
    )


@pytest.mark.parametrize("directive", _PARTIAL_FULL)
def test_svg_embeds_without_crash(directive, tmp_path):
    """Native SVG for each PARTIAL/FULL type must embed into a PPTX slide without crashing."""
    src = _PPTX_COMPAT_SOURCES.get(directive)
    if src is None:
        pytest.skip(f"no minimal source defined for {directive!r}")

    svg_str = dispatch_native(src)
    assert svg_str and "<svg" in svg_str, f"{directive}: dispatch_native returned no SVG"

    svg_file = tmp_path / f"{directive}.svg"
    svg_file.write_text(svg_str, encoding="utf-8")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    conv = SvgConverter()
    conv.convert(svg_file, slide)

    assert conv.stats["errors"] == 0, (
        f"{directive}: SvgConverter reported {conv.stats['errors']} error(s); "
        f"stats={conv.stats}"
    )
    assert conv.stats.get("shapes", 0) > 0, (
        f"{directive}: SvgConverter embedded 0 shapes — "
        f"embed path did not exercise the SVG content; stats={conv.stats}"
    )
