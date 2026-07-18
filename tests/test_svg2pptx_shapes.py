"""SVG basic-shape coverage for scripts/svg2pptx.py.

Import pattern mirrors tests/test_mermaid_layout.py: sys.path.insert so tests/
doesn't need a conftest.py. Drives SvgConverter.convert on in-memory probe SVGs
and asserts on the observable result (the stats counters and emitted shape XML).

Regression guard for the silently-dropped <polygon> arrowheads: the mermaid
renderer emits every arrowhead as a <polygon>, so a drop here loses all arrows
in the PPTX export.
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_ROOT / "scripts"))

import svg2pptx  # noqa: E402
from svg2pptx import SvgConverter  # noqa: E402

SVG_HEAD = ('<svg xmlns="http://www.w3.org/2000/svg" width="400" height="300" '
            'viewBox="0 0 400 300">')


def _convert(svg_body: str, tmp_path: Path):
    """Write an SVG, convert it onto a fresh blank slide, return (stats, slide_xml)."""
    f = tmp_path / "probe.svg"
    f.write_text(f"{SVG_HEAD}{svg_body}</svg>")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    conv = SvgConverter()
    conv.convert(f, slide)
    xml = etree.tostring(slide._element).decode()
    return conv.stats, xml


def test_all_four_basic_shapes_convert(tmp_path):
    """rect + path + polygon + polyline all become native shapes (was 2)."""
    stats, _ = _convert(
        '<rect x="10" y="10" width="80" height="40" fill="#333"/>'
        '<path d="M 200 100 L 192 96 L 192 104 Z" fill="#e00"/>'
        '<polygon points="300,100 292,96 292,104" fill="#0a0"/>'
        '<polyline points="20,200 120,200 120,260" fill="none" stroke="#00f" stroke-width="2"/>',
        tmp_path,
    )
    assert stats["shapes"] == 4, stats
    assert stats["errors"] == 0, stats


def test_polygon_preserves_fill(tmp_path):
    """A filled polygon carries its fill colour through to the shape XML."""
    _, xml = _convert('<polygon points="300,100 292,96 292,104" fill="#00AA00"/>', tmp_path)
    assert "00AA00" in xml.upper()


def test_polygon_is_closed(tmp_path):
    """polygon geometry must close (emits an a:close), unlike polyline."""
    _, xml = _convert('<polygon points="300,100 292,96 292,104" fill="#0a0"/>', tmp_path)
    assert "a:close" in xml


def test_polyline_is_open(tmp_path):
    """polyline must not emit a closing segment (no a:close in its geometry)."""
    _, xml = _convert(
        '<polyline points="20,200 120,200 120,260" fill="none" stroke="#00f" stroke-width="2"/>',
        tmp_path,
    )
    assert "a:close" not in xml


def test_arrowhead_polygon_converts(tmp_path):
    """The renderer's actual _arrowhead point-string converts (the real bug)."""
    from mermaid_layout._routing import _arrowhead
    pts = _arrowhead(300, 100, -1, 0)  # leftward tip, as LR edges emit
    stats, _ = _convert(f'<polygon points="{pts}" fill="#64748b"/>', tmp_path)
    assert stats["shapes"] == 1, stats


def test_unhandled_leaf_is_counted_not_swallowed(tmp_path):
    """An unknown leaf rendering element surfaces via the unhandled counter."""
    stats, _ = _convert('<foobar x="1" y="2"/>', tmp_path)
    assert stats.get("unhandled", 0) >= 1, stats
    assert stats["errors"] == 0, stats


def test_comment_nodes_are_not_counted_unhandled(tmp_path):
    """dom-to-svg emits XML comments (falsy tag) — they must not count as unhandled."""
    stats, _ = _convert(
        '<!-- a comment --><rect x="10" y="10" width="80" height="40" fill="#333"/>',
        tmp_path,
    )
    assert stats.get("unhandled", 0) == 0, stats
    assert stats["shapes"] == 1, stats


def test_metadata_leaves_are_exempt(tmp_path):
    """Non-rendering leaves (title/desc/metadata) do not count as unhandled."""
    stats, _ = _convert(
        '<title>t</title><desc>d</desc><rect x="10" y="10" width="80" height="40" fill="#333"/>',
        tmp_path,
    )
    assert stats.get("unhandled", 0) == 0, stats
    assert stats["shapes"] == 1, stats


# ── <marker> / marker-end resolution (edge arrowheads after renderer #55) ──────

# The renderer's arrow-normal marker, verbatim from mermaid_layout._renderer.
_MARKER = (
    '<defs><marker id="ah" viewBox="0 -4 9 8" refX="9" refY="0" markerWidth="9" '
    'markerHeight="8" markerUnits="userSpaceOnUse" orient="auto">'
    '<polygon points="0,-4 9,0 0,4" fill="#334455"/></marker></defs>'
)


_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _offsets(xml):
    """Every shape's top-left (a:off) in px, back-converted from EMU."""
    from svg2pptx import EMU_PX
    root = etree.fromstring(xml.encode())
    return [(int(o.get("x")) / EMU_PX, int(o.get("y")) / EMU_PX)
            for o in root.iter(f"{_A}off")]


def _boxes(xml):
    """Every shape's (x, y, w, h) in px — pairs each a:off with its sibling a:ext."""
    from svg2pptx import EMU_PX
    root = etree.fromstring(xml.encode())
    out = []
    for xfrm in root.iter(f"{_A}xfrm"):
        off, ext = xfrm.find(f"{_A}off"), xfrm.find(f"{_A}ext")
        if off is not None and ext is not None:
            out.append((int(off.get("x")) / EMU_PX, int(off.get("y")) / EMU_PX,
                        int(ext.get("cx")) / EMU_PX, int(ext.get("cy")) / EMU_PX))
    return out


def _near(offs, x, y, tol=2):
    return any(abs(px - x) <= tol and abs(py - y) <= tol for px, py in offs)


def test_marker_end_adds_arrowhead_shape(tmp_path):
    """A path with marker-end draws the marker geometry as an extra shape."""
    without, _ = _convert(
        '<path d="M 100 100 L 100 300" stroke="#334455" stroke-width="2" fill="none"/>',
        tmp_path,
    )
    with_marker, _ = _convert(
        _MARKER + '<path d="M 100 100 L 100 300" stroke="#334455" stroke-width="2" '
        'fill="none" marker-end="url(#ah)"/>',
        tmp_path,
    )
    assert with_marker["shapes"] == without["shapes"] + 1, (without, with_marker)
    assert with_marker["errors"] == 0


def test_marker_end_oriented_at_endpoint(tmp_path):
    """orient=auto: a downward path puts the arrowhead just above its endpoint (100,300).

    Transform of the arrow-normal triangle rotated to point down gives a bbox of
    roughly x∈[96,104], y∈[291,300] — top-left ≈ (96, 291).
    """
    _, xml = _convert(
        _MARKER + '<path d="M 100 100 L 100 300" stroke="#334455" stroke-width="2" '
        'fill="none" marker-end="url(#ah)"/>',
        tmp_path,
    )
    assert _near(_offsets(xml), 96, 291), _offsets(xml)


def test_marker_orientation_differs_by_direction(tmp_path):
    """A rightward path orients the arrowhead differently from a downward one."""
    _, down = _convert(
        _MARKER + '<path d="M 100 100 L 100 300" fill="none" stroke="#334455" marker-end="url(#ah)"/>',
        tmp_path,
    )
    _, right = _convert(
        _MARKER + '<path d="M 100 100 L 300 100" fill="none" stroke="#334455" marker-end="url(#ah)"/>',
        tmp_path,
    )
    # Downward arrow tip at (100,300), top-left ≈ (96,291); rightward at (300,100) ≈ (291,96).
    assert _near(_offsets(down), 96, 291), _offsets(down)
    assert _near(_offsets(right), 291, 96), _offsets(right)


_START_MARKER = _MARKER.replace('id="ah"', 'id="ah2"').replace(
    'orient="auto"', 'orient="auto-start-reverse"')


def test_marker_start_auto_reverse(tmp_path):
    """auto-start-reverse flips a start marker to point back along the path.

    Downward path start (100,100): plain auto would point down (base above,
    top-left y≈91); reversed points up (base below, top-left ≈ (96,100)).
    """
    _, xml = _convert(
        _START_MARKER + '<path d="M 100 100 L 100 300" fill="none" stroke="#334455" '
        'marker-start="url(#ah2)"/>',
        tmp_path,
    )
    assert _near(_offsets(xml), 96, 100), _offsets(xml)


def test_marker_fixed_numeric_orient(tmp_path):
    """A numeric orient ignores the tangent and rotates by that many degrees."""
    fixed = _MARKER.replace('id="ah"', 'id="ah3"').replace('orient="auto"', 'orient="90"')
    _, xml = _convert(
        fixed + '<path d="M 100 100 L 300 100" fill="none" stroke="#334455" '
        'marker-end="url(#ah3)"/>',
        tmp_path,
    )
    # Rightward path (auto would give top-left ≈ (291,96)); fixed 90° gives ≈ (296,91).
    assert _near(_offsets(xml), 296, 91), _offsets(xml)


def test_marker_stroke_width_units_scale(tmp_path):
    """markerUnits=strokeWidth scales the arrowhead by the host stroke-width."""
    sw = (_MARKER.replace('id="ah"', 'id="ah4"')
          .replace('markerUnits="userSpaceOnUse"', 'markerUnits="strokeWidth"'))
    _, xml = _convert(
        sw + '<path d="M 100 100 L 100 300" fill="none" stroke="#334455" stroke-width="2" '
        'marker-end="url(#ah4)"/>',
        tmp_path,
    )
    # The base marker is 8×9 px; stroke-width 2 doubles it to ~16×18.
    arrow = [b for b in _boxes(xml) if abs(b[2] - 16) <= 2 and abs(b[3] - 18) <= 2]
    assert arrow, _boxes(xml)


def test_marker_path_child_renders(tmp_path):
    """A <path>-child marker (the renderer's arrow-open) exercises _transform_path_d."""
    open_marker = (
        '<defs><marker id="ao" viewBox="0 -4 9 8" refX="9" refY="0" markerWidth="9" '
        'markerHeight="8" markerUnits="userSpaceOnUse" orient="auto">'
        '<path d="M 0,-4 L 9,0 L 0,4" fill="none" stroke="#e8924a" stroke-width="1.5"/>'
        '</marker></defs>'
    )
    stats, xml = _convert(
        open_marker + '<path d="M 100 100 L 100 300" fill="none" stroke="#e8924a" '
        'marker-end="url(#ao)"/>',
        tmp_path,
    )
    assert stats["shapes"] == 2, stats  # edge line + arrowhead
    assert _near(_offsets(xml), 96, 291), _offsets(xml)


def test_marker_on_relative_path_is_skipped(tmp_path):
    """A relative/H/V host path can't be paired safely, so the marker is skipped
    (base shape still converts, no crash)."""
    stats, _ = _convert(
        _MARKER + '<path d="M 100 100 h 50 v 50" fill="none" stroke="#334455" '
        'marker-end="url(#ah)"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats  # only the base path; marker safely skipped
    assert stats["errors"] == 0, stats


def test_orient_radians_units():
    """Unitless orient is degrees; grad is matched before its rad suffix."""
    import math
    f = SvgConverter._orient_radians
    assert f("90") == pytest.approx(math.pi / 2)      # unitless -> degrees
    assert f("90deg") == pytest.approx(math.pi / 2)
    assert f("1rad") == pytest.approx(1.0)
    assert f("200grad") == pytest.approx(math.pi)     # not mis-read as radians
    assert f("0.5turn") == pytest.approx(math.pi)
    assert f("junk") == 0.0


def test_marker_malformed_viewbox_no_crash(tmp_path):
    """A <4-number viewBox degrades gracefully instead of aborting the conversion."""
    bad = _MARKER.replace('viewBox="0 -4 9 8"', 'viewBox="0 -4 9"')
    stats, _ = _convert(
        bad + '<path d="M 100 100 L 100 300" fill="none" stroke="#334455" '
        'marker-end="url(#ah)"/>',
        tmp_path,
    )
    assert stats["errors"] == 0, stats
    assert stats["shapes"] >= 1, stats


def test_unreferenced_defs_marker_not_drawn(tmp_path):
    """A marker defined but never referenced produces no shape (defs is skipped)."""
    stats, _ = _convert(_MARKER, tmp_path)
    assert stats["shapes"] == 0, stats
    assert stats.get("unhandled", 0) == 0, stats


def test_marker_missing_ref_does_not_crash(tmp_path):
    """marker-end to a missing id: base shape still converts, no error."""
    stats, _ = _convert(
        '<path d="M 100 100 L 100 300" stroke="#334455" fill="none" marker-end="url(#nope)"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert stats["errors"] == 0, stats


# ── straight H/V/L connector paths (the missing-connector-lines bug) ───────────
#
# Mermaid/dom-to-svg emit orthogonal edges as single- or multi-segment paths
# built from H (horizontal) and V (vertical) commands, e.g. "M118 67 H153".
# Two defects dropped or mis-placed them in the PPTX export:
#   (1) the len(nums) < 4 guard dropped any single-segment H/V path (3 numbers);
#   (2) the bbox was computed by naive number-pairing (xs=coords[0::2]), which
#       ignores H/V current-point semantics and transposes multi-segment boxes.

def test_pure_horizontal_connector_converts(tmp_path):
    """A single-segment horizontal connector (M x y H x2 -> 3 numbers) renders.

    Was dropped by `if len(nums) < 4: return`, taking most of a diagram's
    orthogonal edges with it.
    """
    stats, _ = _convert(
        '<path d="M118 67 H153" fill="none" stroke="rgb(63,125,90)" stroke-width="1.7px"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert stats["errors"] == 0, stats


def test_pure_vertical_connector_converts(tmp_path):
    """A single-segment vertical connector (M x y V y2 -> 3 numbers) renders."""
    stats, _ = _convert(
        '<path d="M241 110 V125" fill="none" stroke="rgb(107,74,122)" stroke-width="1.7px"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert stats["errors"] == 0, stats


def test_horizontal_connector_carries_stroke(tmp_path):
    """The connector must render as a visible line — its stroke colour survives."""
    _, xml = _convert(
        '<path d="M118 67 H153" fill="none" stroke="rgb(63,125,90)" stroke-width="1.7px"/>',
        tmp_path,
    )
    assert "3F7D5A".upper() in xml.upper(), "connector stroke colour rgb(63,125,90) missing"


def test_horizontal_connector_bbox(tmp_path):
    """A pure-H connector's box spans the segment (x∈[118,153]) and is ~flat."""
    _, xml = _convert(
        '<path d="M118 67 H153" fill="none" stroke="rgb(63,125,90)" stroke-width="1.7px"/>',
        tmp_path,
    )
    boxes = _boxes(xml)
    assert any(abs(x - 118) <= 2 and abs(y - 67) <= 2 and abs(w - 35) <= 2 and h <= 2
               for x, y, w, h in boxes), boxes


def test_multisegment_connector_bbox_from_traced_points(tmp_path):
    """A V-then-H-then-V connector gets its true traced bbox, not a transposed one.

    'M331 70 V90 H586 V79' visits (331,70),(331,90),(586,90),(586,79):
    true bbox = (x=331, y=70, w=255, h=20). The old naive number-pairing gave a
    wildly wrong (x=79, y=70, w=252, h=516) — the line landed off-place and huge.
    """
    _, xml = _convert(
        '<path d="M331 70 V90 H586 V79" fill="none" stroke="rgb(31,58,95)" stroke-width="1.7px"/>',
        tmp_path,
    )
    boxes = _boxes(xml)
    assert any(abs(x - 331) <= 2 and abs(y - 70) <= 2 and abs(w - 255) <= 2 and abs(h - 20) <= 2
               for x, y, w, h in boxes), boxes
    # And explicitly reject the transposed naive bbox.
    assert not any(abs(x - 79) <= 2 and abs(h - 516) <= 4 for x, y, w, h in boxes), boxes


# ── icon-detail paths (regression guard for the bbox rewrite) ──────────────────
#
# Correcting the connector bbox also corrected it for small icon sub-paths, which
# the old naive number-pairing had accidentally inflated (relative deltas read as
# absolute coords) so they slipped past the tiny-path filter. Two things must
# hold so icons don't lose strokes: (a) the bbox must include curve/arc
# endpoints, not just the M/L/H/V vertices; (b) a tiny *stroked* path (icon
# detail, a tick, a chevron) must still render — only tiny *fill-only*
# decoration is dropped.

def test_tiny_stroked_chevron_renders(tmp_path):
    """A ~2×3px stroked chevron (real Lucide icon detail) still renders as a line."""
    stats, _ = _convert(
        '<path d="M9 13l-2 1.5L9 16" fill="none" stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert stats["errors"] == 0, stats


def test_tiny_stroked_chevron_bbox_is_correct(tmp_path):
    """The chevron sits at its true position/size (≈7,13 / 2×3), not an inflated box."""
    _, xml = _convert(
        '<path d="M9 13l-2 1.5L9 16" fill="none" stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    boxes = _boxes(xml)
    assert any(abs(x - 7) <= 1 and abs(y - 13) <= 1 and abs(w - 2) <= 1 and abs(h - 3) <= 1
               for x, y, w, h in boxes), boxes


def test_icon_path_mixing_hv_and_arcs_renders(tmp_path):
    """A stroked icon path mixing h/v with arcs still renders (regression guard).

    Arcs now render as béziers and the bbox spans the arc envelope too, so this
    small stroked sub-path must not be dropped and must not crash.
    """
    stats, _ = _convert(
        '<path d="M18 10h2a1 1 0 0 1 1 1v2a1 1 0 0 1-1 1h-2" fill="none" '
        'stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert stats["errors"] == 0, stats


def _path_pts_range(xml):
    """(x-range, y-range) of the emitted a:pathLst point coords, in OOXML units."""
    root = etree.fromstring(xml.encode())
    xs = [int(p.get("x")) for p in root.iter(f"{_A}pt")]
    ys = [int(p.get("y")) for p in root.iter(f"{_A}pt")]
    return (min(xs), max(xs)), (min(ys), max(ys))


def test_drawn_geometry_fills_its_bbox(tmp_path):
    """The custGeom points must fill 0..100000 in both axes — i.e. the drawn
    geometry is consistent with the shape's box.

    This is the invariant the two bugs violated: the naive-number-pairing bbox
    (bug 2) put the box in the wrong place, so the same points normalized outside
    0..100000 / squished into a corner. A single walker for both box and geometry
    keeps them in lockstep.
    """
    _, xml = _convert(
        '<path d="M331 70 V90 H586 V79" fill="none" stroke="rgb(31,58,95)" stroke-width="1.7px"/>',
        tmp_path,
    )
    (xlo, xhi), (ylo, yhi) = _path_pts_range(xml)
    assert xlo == 0 and xhi == 100000, (xlo, xhi)   # spans the box, not beyond
    assert ylo == 0 and yhi == 100000, (ylo, yhi)


def test_tiny_fill_only_decoration_still_dropped(tmp_path):
    """A tiny fill-only path (no stroke) is still filtered as decorative noise."""
    stats, _ = _convert(
        '<path d="M10 10 L11 10 L11 11 Z" fill="#333"/>',
        tmp_path,
    )
    assert stats["shapes"] == 0, stats
    assert stats["skipped"] >= 1, stats


def test_compound_path_geometry_fills_its_bbox(tmp_path):
    """Z-then-relative compound path: box and geometry stay in lockstep.

    `_path_points` resets the pen to the subpath start on Z; the geometry builder
    must do the same, or the second (relative) subpath is drawn from the wrong
    origin and its points normalize outside 0..100000 (the bbox it claims to fill).
    """
    _, xml = _convert(
        '<path d="M0 0 H10 V10 Z m0 20 h10" fill="none" stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    (xlo, xhi), (ylo, yhi) = _path_pts_range(xml)
    assert 0 <= xlo and xhi <= 100000, (xlo, xhi)
    assert 0 <= ylo and yhi <= 100000, (ylo, yhi)
    assert xhi == 100000 and yhi == 100000, ((xlo, xhi), (ylo, yhi))


# ── line thickness, solid vs dashed, arrowheads ───────────────────────────────

def _ln_width(xml):
    """The first a:ln @w (EMU) in the slide, or None."""
    root = etree.fromstring(xml.encode())
    for ln in root.iter(f"{_A}ln"):
        if ln.get("w"):
            return int(ln.get("w"))
    return None


def test_connector_line_thickness_matches_pixels(tmp_path):
    """A 1.7px stroke renders 1.7 slide-px wide (EMU_PX=9525), not 1.7pt (12700).

    Geometry is placed at 9525 EMU/px; line weight must use the same scale or
    every connector reads ~1.33× too heavy against the HTML.
    """
    from svg2pptx import EMU_PX
    _, xml = _convert(
        '<path d="M118 67 H153" fill="none" stroke="#3f7d5a" stroke-width="1.7px"/>',
        tmp_path,
    )
    w = _ln_width(xml)
    assert w == round(1.7 * EMU_PX), (w, round(1.7 * EMU_PX))


def test_dashed_connector_emits_prstdash(tmp_path):
    """A dashed connector carries a dash style, not a solid line."""
    _, xml = _convert(
        '<path d="M320 67 H635" fill="none" stroke="#6b4a7a" stroke-width="1.7px" '
        'stroke-dasharray="5px, 4px"/>',
        tmp_path,
    )
    assert "prstDash" in xml, "dashed connector rendered as solid (no a:prstDash)"


def test_dotted_connector_maps_to_dot(tmp_path):
    """A short-dash array (dots) maps to a dotted preset, a long one to dash."""
    _, dotted = _convert(
        '<path d="M694 194 H375" fill="none" stroke="#3f7d5a" stroke-width="1.9px" '
        'stroke-dasharray="2px, 4px"/>',
        tmp_path,
    )
    _, dashed = _convert(
        '<path d="M320 67 H635" fill="none" stroke="#6b4a7a" stroke-width="1.7px" '
        'stroke-dasharray="5px, 4px"/>',
        tmp_path,
    )
    assert 'val="sysDot"' in dotted or 'val="dot"' in dotted, dotted[dotted.find("prstDash")-5:dotted.find("prstDash")+40] if "prstDash" in dotted else "no prstDash"
    assert 'val="dash"' in dashed


def test_solid_connector_has_no_dash(tmp_path):
    """A connector without stroke-dasharray stays solid (no dash element)."""
    _, xml = _convert(
        '<path d="M118 67 H153" fill="none" stroke="#3f7d5a" stroke-width="1.7px"/>',
        tmp_path,
    )
    assert "prstDash" not in xml and "custDash" not in xml


def test_hollow_circle_stroke_thickness_matches_pixels(tmp_path):
    """A fill=none stroked circle's outline uses the px scale too (not 1pt)."""
    from svg2pptx import EMU_PX
    _, xml = _convert(
        '<circle cx="50" cy="50" r="20" fill="none" stroke="#333" stroke-width="2"/>',
        tmp_path,
    )
    assert _ln_width(xml) == round(2 * EMU_PX), _ln_width(xml)


def test_dash_preset_handles_bad_and_zero_arrays():
    """dash_preset degrades to solid on non-numeric / zero arrays (no crash)."""
    from svg2pptx import dash_preset
    assert dash_preset("inherit") is None       # unparseable -> solid, not ValueError
    assert dash_preset("0") is None             # zero-length dash is solid in SVG
    assert dash_preset("0 0") is None
    assert dash_preset("") is None
    assert dash_preset("none") is None
    assert dash_preset("5px, 4px", 1.7) == "dash"
    assert dash_preset("2px, 4px", 1.9) == "sysDot"


def test_connector_with_unparseable_dasharray_still_renders(tmp_path):
    """A connector with a bad stroke-dasharray falls back to solid, not dropped."""
    stats, xml = _convert(
        '<path d="M118 67 H153" fill="none" stroke="#3f7d5a" stroke-width="1.7px" '
        'stroke-dasharray="inherit"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert stats["errors"] == 0, stats
    assert "prstDash" not in xml  # solid fallback


def test_inline_continuation_tspan_keeps_leading_space(tmp_path):
    """A bold label followed by a positioned ' — rest' tspan keeps its leading
    space, else 'TEMPORARY — Atlas' collapses to 'TEMPORARY—Atlas'."""
    _, xml = _convert(
        '<text fill="#333" font-size="14px">'
        '<tspan x="10" y="20">TEMPORARY</tspan>'
        '<tspan x="100" y="20"> — Atlas reads</tspan></text>',
        tmp_path,
    )
    root = etree.fromstring(xml.encode())
    texts = [t.text for t in root.iter(f"{_A}t")]
    assert any(t and t.startswith(" ") for t in texts), texts


def test_horizontal_line_is_flat_not_slanted(tmp_path):
    """A horizontal <line> renders flat (0 minor dimension), not slanted.

    The `line` preset draws corner-to-corner of its box; the old `or 1` forced a
    1px minor dimension, so a 20px horizontal icon line got a 1px vertical rise
    (~a few degrees of slant, very visible on short icon lines)."""
    _, xml = _convert(
        '<line x1="2" y1="7" x2="22" y2="7" stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    boxes = _boxes(xml)
    assert any(abs(w - 20) <= 1 and h <= 0.2 for x, y, w, h in boxes), boxes


def test_vertical_line_is_flat_not_slanted(tmp_path):
    """A vertical <line> renders flat (0 minor dimension), not slanted."""
    _, xml = _convert(
        '<line x1="7" y1="7" x2="7" y2="21" stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    boxes = _boxes(xml)
    assert any(w <= 0.2 and abs(h - 14) <= 1 for x, y, w, h in boxes), boxes


def test_diagonal_line_keeps_its_slope(tmp_path):
    """A genuinely diagonal <line> keeps both dimensions (not flattened)."""
    _, xml = _convert(
        '<line x1="2" y1="2" x2="20" y2="14" stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    boxes = _boxes(xml)
    assert any(abs(w - 18) <= 1 and abs(h - 12) <= 1 for x, y, w, h in boxes), boxes


def test_round_linecap_emitted(tmp_path):
    """stroke-linecap=round -> a:ln cap="rnd" (makes dot-grid icons visible)."""
    _, xml = _convert(
        '<path d="M8 8h.01M12 8h.01M16 8h.01" fill="none" stroke="#333" '
        'stroke-width="1.5" stroke-linecap="round"/>',
        tmp_path,
    )
    assert 'cap="rnd"' in xml, "round linecap not emitted"


def test_round_linejoin_emitted(tmp_path):
    """stroke-linejoin=round -> an a:round child in the line (Lucide corners)."""
    _, xml = _convert(
        '<path d="M2 2 L20 2 L20 20" fill="none" stroke="#333" stroke-width="1.5" '
        'stroke-linejoin="round"/>',
        tmp_path,
    )
    assert "a:round" in xml, "round linejoin not emitted"


def test_default_stroke_has_no_cap_or_join(tmp_path):
    """A plain connector (no linecap/linejoin) stays flat/miter — no cap attr."""
    _, xml = _convert(
        '<path d="M118 67 H153" fill="none" stroke="#3f7d5a" stroke-width="1.7px"/>',
        tmp_path,
    )
    assert 'cap="rnd"' not in xml and 'cap="sq"' not in xml


def test_dot_grid_icon_dots_render(tmp_path):
    """The vector-database dot grid (9 near-zero `h.01` segments + round caps)
    renders as a shape (was invisible with butt caps)."""
    stats, _ = _convert(
        '<path d="M8 8h.01M12 8h.01M16 8h.01M8 12h.01M12 12h.01M16 12h.01'
        'M8 16h.01M12 16h.01M16 16h.01" fill="none" stroke="#333" '
        'stroke-width="1.5" stroke-linecap="round"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert stats["errors"] == 0, stats


def test_path_in_scaled_group_positions_correctly(tmp_path):
    """A path inside a scaled <g> places at local*scale + accumulated offset.

    Guards the icon-scattering fix (`bx*scale+ox`, not `(bx+ox)*scale`): all other
    tests run at scale 1 where the two are identical, so this is the only test
    that exercises the scaled-<g> transform the icons actually use.
    """
    from svg2pptx import EMU_PX
    _, xml = _convert(
        '<g transform="matrix(1.08333 0 0 1.08333 450 503)">'
        '<path d="M9 13 L11 16" fill="none" stroke="#333" stroke-width="1.5"/></g>',
        tmp_path,
    )
    # origin should be bx*scale + translate = 9*1.08333 + 450 ≈ 459.75, 13*1.08333+503 ≈ 517.08
    boxes = _boxes(xml)
    assert any(abs(x - (9 * 1.08333 + 450)) <= 1 and abs(y - (13 * 1.08333 + 503)) <= 1
               for x, y, w, h in boxes), boxes
    # and NOT the double-scaled origin (9+450)*1.08333 ≈ 497.25
    assert not any(abs(x - (9 + 450) * 1.08333) <= 1 for x, y, w, h in boxes), boxes


def test_arc_command_renders_as_bezier(tmp_path):
    """An elliptical-arc (`A`) segment is drawn as cubic béziers, not skipped.

    Lucide icons build rounded corners and circular parts from `a` commands; when
    the builder skipped them the icons rendered as fragments (people icon -> a dot,
    rounded cards -> square brackets)."""
    stats, xml = _convert(
        '<path d="M0 10 A10 10 0 0 1 20 10" fill="none" stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert "cubicBezTo" in xml, "arc was not converted to bézier geometry"


def test_arc_spans_its_full_extent(tmp_path):
    """A semicircle arc from (0,10) to (20,10) bounds the full ~20-wide span."""
    _, xml = _convert(
        '<path d="M0 10 A10 10 0 0 1 20 10" fill="none" stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    boxes = _boxes(xml)
    assert any(abs(x - 0) <= 2 and abs(w - 20) <= 2 for x, y, w, h in boxes), boxes


def test_relative_arc_icon_path_renders(tmp_path):
    """The deck's people-icon arc body (relative `a`) renders as a shape."""
    stats, xml = _convert(
        '<path d="M3.7 19a5.3 5.3 0 0 1 10.6 0" fill="none" stroke="#3f7d5a" stroke-width="1.5"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert "cubicBezTo" in xml
    boxes = _boxes(xml)
    assert any(abs(w - 10.6) <= 2 for x, y, w, h in boxes), boxes  # spans the arc width


def test_arc_geometry_fills_its_bbox(tmp_path):
    """Arc bbox (walker) and drawn béziers (builder) stay in lockstep."""
    _, xml = _convert(
        '<path d="M0 10 A10 10 0 0 1 20 10" fill="none" stroke="#333" stroke-width="1.5"/>',
        tmp_path,
    )
    (xlo, xhi), (ylo, yhi) = _path_pts_range(xml)
    assert xlo == 0 and xhi == 100000, (xlo, xhi)
    assert ylo == 0 and yhi == 100000, (ylo, yhi)


def test_polygon_arrowhead_renders_at_line_end(tmp_path):
    """The deck draws arrowheads as standalone filled polygons at the edge tip;
    they must render as a filled triangle shape (arrows are not <marker>s here)."""
    stats, xml = _convert(
        '<polygon points="529,246.0 538,251.0 529,256.0" fill="rgb(63,125,90)" stroke="none"/>',
        tmp_path,
    )
    assert stats["shapes"] == 1, stats
    assert "3F7D5A" in xml.upper()  # arrowhead fill colour preserved
    # tip sits at (538,251); base spans x∈[529,538], y∈[246,256]
    assert any(abs(x - 529) <= 2 and abs(y - 246) <= 2 for x, y, w, h in _boxes(xml)), _boxes(xml)
